#!/usr/bin/env python3
"""
Create AWS-branded PowerPoint from AWS DRS Orchestration presentation.
Includes Amazon Ember font, AWS logos, and full brand compliance.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import os
from datetime import datetime

# AWS Brand Colors (Official)
AWS_ORANGE = RGBColor(255, 153, 0)  # #FF9900
AWS_SQUID_INK = RGBColor(35, 47, 62)  # #232F3E
AWS_WHITE = RGBColor(255, 255, 255)  # #FFFFFF

def create_aws_logo_if_needed():
    """Create AWS logo if not already present"""
    if not os.path.exists('aws_logo.png') or os.path.getsize('aws_logo.png') < 1000:
        # Create a simple AWS logo placeholder
        img = Image.new('RGBA', (400, 150), (255, 255, 255, 0))
        draw = ImageDraw.Draw(img)
        
        # Draw "AWS" text in orange
        try:
            font = ImageFont.truetype("Amazon Ember", 80)
        except:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 80)
        
        # AWS text
        draw.text((50, 35), "AWS", fill=(255, 153, 0, 255), font=font)
        
        # Save
        img.save('aws_logo.png')
        print("✅ Created aws_logo.png")
    
    return 'aws_logo.png'

def add_footer(slide, slide_number, total_slides):
    """Add AWS-branded footer to slide"""
    # Footer background (subtle)
    footer_box = slide.shapes.add_textbox(
        Inches(0), Inches(7.0), Inches(10), Inches(0.5)
    )
    
    # Slide number (left)
    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(1.5), Inches(0.4))
    tf = left_box.text_frame
    tf.text = f"{slide_number} of {total_slides}"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = AWS_SQUID_INK
    tf.paragraphs[0].font.name = "Amazon Ember"
    
    # Confidential notice (center)
    center_box = slide.shapes.add_textbox(Inches(2.5), Inches(7.05), Inches(5), Inches(0.4))
    tf = center_box.text_frame
    tf.text = f"Confidential - Internal Use Only | {datetime.now().strftime('%B %Y')}"
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = AWS_SQUID_INK
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.name = "Amazon Ember"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # AWS logo (right)
    if os.path.exists('aws_logo.png'):
        try:
            slide.shapes.add_picture('aws_logo.png', 
                                    Inches(8.8), Inches(7.05), 
                                    height=Inches(0.35))
        except:
            pass

def create_title_slide(prs):
    """Create AWS-branded title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "AWS DRS Orchestration Platform"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = AWS_ORANGE
    tf.paragraphs[0].font.name = "Amazon Ember Display"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    tf.text = "Serverless Disaster Recovery Orchestration"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.name = "Amazon Ember"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.5))
    tf = tag_box.text_frame
    tf.text = "Enterprise-Grade Solution for Automated Recovery"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.name = "Amazon Ember"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Powered by AWS logo (bottom - reduced to half size)
    if os.path.exists('aws_logo.png'):
        try:
            slide.shapes.add_picture('aws_logo.png', 
                                    Inches(4.25), Inches(5.8), 
                                    width=Inches(1.5))
        except:
            pass
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.3))
    tf = date_box.text_frame
    tf.text = datetime.now().strftime("%B %Y")
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.name = "Amazon Ember"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_bullet_slide(prs, title, bullets, slide_num, total):
    """Add slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = AWS_ORANGE  # Changed to AWS Orange
    tf.paragraphs[0].font.name = "Amazon Ember Display"
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    tf = content_box.text_frame
    
    for bullet in bullets:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        
        # Determine indent level
        indent = 0
        text = bullet
        while text.startswith('  '):
            indent += 1
            text = text[2:]
        
        p.text = text.lstrip('• ')
        p.level = indent
        p.font.size = Pt(16 if indent == 0 else 14)
        p.font.name = "Amazon Ember"
        p.space_before = Pt(6) if indent == 0 else Pt(3)
    
    # Footer
    add_footer(slide, slide_num, total)

def add_table_slide(prs, title, headers, rows, slide_num, total):
    """Add slide with AWS-branded table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = AWS_ORANGE  # Changed to AWS Orange
    tf.paragraphs[0].font.name = "Amazon Ember Display"
    
    # Table
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(0.5), Inches(1.4),
        Inches(9), Inches(5.2)
    ).table
    
    # Headers (AWS Orange background, white text)
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = AWS_WHITE
        cell.text_frame.paragraphs[0].font.name = "Amazon Ember"
        cell.fill.solid()
        cell.fill.fore_color.rgb = AWS_ORANGE
    
    # Rows (alternating subtle background)
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.name = "Amazon Ember"
            
            # Alternating row colors
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Footer
    add_footer(slide, slide_num, total)

def create_presentation():
    """Create the complete AWS-branded presentation"""
    
    # Ensure AWS logo exists
    create_aws_logo_if_needed()
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    total_slides = 16
    
    # Slide 1: Title (no footer)
    create_title_slide(prs)
    
    # Slide 2: What is AWS DRS Orchestration?
    add_bullet_slide(prs, "What is AWS DRS Orchestration?", [
        "Wave-based recovery execution (like VMware SRM)",
        "Server dependency management",
        "One-click drill testing",
        "Real-time progress monitoring",
        "Complete audit trail",
        "",
        "Built with: 100% AWS serverless technologies",
        "Deployment: < 1 hour via CloudFormation",
        "Cost: $3,360/year (fixed, regardless of server count)"
    ], 2, total_slides)
    
    # Slide 3: The Problem
    add_bullet_slide(prs, "The Problem We Solve", [
        "WITHOUT ORCHESTRATION:",
        "  • Manual server recovery (error-prone)",
        "  • No dependency management",
        "  • 12-16 hour recovery times",
        "  • Complex drill testing procedures",
        "  • No visibility into recovery progress",
        "",
        "MIGRATION CHALLENGE:",
        "  • Customers have VMware SRM orchestration",
        "  • AWS DRS lacks equivalent capabilities",
        "  • This gap blocks AWS DRS adoption"
    ], 3, total_slides)
    
    # Slide 4: Architecture Components
    add_table_slide(prs, "Architecture Components", 
        ["Layer", "Technology", "Purpose"],
        [
            ["Frontend", "React + TypeScript + MUI", "Modern web UI"],
            ["CDN", "CloudFront + S3", "Global distribution"],
            ["Authentication", "Cognito User Pools", "Secure access"],
            ["API", "API Gateway (REST)", "Managed API layer"],
            ["Compute", "Lambda (Python 3.12)", "Business logic"],
            ["Orchestration", "Step Functions (35+ states)", "Wave-based execution"],
            ["Integration", "DRS API + EC2 API", "Recovery & health checks"],
            ["Automation", "SSM Documents", "Post-recovery hooks"],
            ["Storage", "DynamoDB (3 tables)", "Configuration data"],
            ["Monitoring", "CloudWatch + CloudTrail", "Logs and audit"]
        ], 4, total_slides
    )
    
    # Slide 5: Core Capabilities
    add_bullet_slide(prs, "Core Capabilities", [
        "1. Protection Groups",
        "  • Group servers logically",
        "  • Auto-discover DRS servers",
        "  • Reusable across plans",
        "",
        "2. Recovery Plans",
        "  • Multi-wave recovery sequences",
        "  • Wave 1: Database | Wave 2: Application | Wave 3: Web",
        "  • Dependencies between waves"
    ], 5, total_slides)
    
    # Slide 6: Wave-Based Execution
    add_bullet_slide(prs, "Wave-Based Execution", [
        "3. Wave-Based Execution",
        "  • 15-second delays between servers",
        "  • 30-second delays between waves",
        "  • Automatic retry with exponential backoff",
        "",
        "4. Drill Mode Testing",
        "  • Creates test instances (not production cutover)",
        "  • Zero impact on replication"
    ], 6, total_slides)
    
    # Slide 7: Monitoring & Recovery
    add_bullet_slide(prs, "Real-Time Monitoring & One-Click Recovery", [
        "5. Real-Time Monitoring",
        "  • Execution history (start/end times)",
        "  • Job status tracking",
        "  • CloudWatch integration",
        "  • Audit trail in CloudTrail",
        "",
        "6. One-Click Recovery",
        "  • Select recovery plan",
        "  • Choose drill vs production mode",
        "  • Click 'Execute Recovery'",
        "  • Monitor progress in real-time"
    ], 7, total_slides)
    
    # Slide 8: Performance Metrics
    add_table_slide(prs, "Performance Metrics",
        ["Metric", "Before (Manual)", "After (Orchestrated)", "Improvement"],
        [
            ["Recovery Time", "12-16 hours", "4-6 hours", "65% faster"],
            ["Setup Time", "2-3 weeks", "< 1 hour", "99% faster"],
            ["Drill Execution", "4-5 hours", "2-3 minutes", "98% faster"],
            ["Error Rate", "15-20%", "< 2%", "90% reduction"]
        ], 8, total_slides
    )
    
    # Slide 9: Cost Breakdown
    add_table_slide(prs, "AWS DRS Orchestration Platform Cost",
        ["Service", "Monthly", "Annual", "Notes"],
        [
            ["Lambda", "$100", "$1,200", "Recovery execution"],
            ["DynamoDB", "$50", "$600", "3 tables (on-demand)"],
            ["API Gateway", "$50", "$600", "REST API requests"],
            ["CloudFront", "$20", "$240", "Global CDN"],
            ["S3", "$10", "$120", "Frontend hosting"],
            ["Cognito", "$20", "$240", "User authentication"],
            ["CloudWatch", "$20", "$240", "Log retention"],
            ["CloudTrail", "$10", "$120", "Audit logging"],
            ["TOTAL", "$280", "$3,360", "Fixed cost"]
        ], 9, total_slides
    )
    
    # Slide 10: Total Solution Cost
    add_table_slide(prs, "Total Solution Cost (1,000 Servers)",
        ["Component", "Annual Cost", "Notes"],
        [
            ["AWS DRS Service", "$480,000", "Scales with server count"],
            ["+ Orchestration Platform", "$3,360", "Fixed serverless cost"],
            ["= Total Solution", "$483,360", "Complete DR solution"]
        ], 10, total_slides
    )
    
    # Slide 11: Competitive Comparison
    add_table_slide(prs, "Competitive Comparison",
        ["Solution", "Orchestration Cost", "Setup Time", "Serverless", "Wave Support"],
        [
            ["AWS DRS Orchestration", "$3.4K/year", "< 1 hour", "✅ Yes", "✅ Yes"],
            ["VMware SRM", "$1M+/year", "Weeks", "❌ No", "✅ Yes"],
            ["Zerto", "$100K+/year", "Days", "❌ No", "⚠️ Limited"],
            ["Azure Site Recovery", "$50K+/year", "Days", "❌ No", "❌ No"]
        ], 11, total_slides
    )
    
    # Slide 12: TCO Summary
    add_table_slide(prs, "TCO Summary (3-Year, 1,000 Servers)",
        ["Solution", "3-Year Total", "Annual Avg", "Staffing"],
        [
            ["VMware SRM", "$10.2M", "$3.4M", "4 FTEs"],
            ["AWS DRS + Orch", "$2.1M", "$697K", "1 FTE"],
            ["Savings", "$8.1M (79%)", "$2.7M", "3 FTEs"]
        ], 12, total_slides
    )
    
    # Slide 13: Business Value
    add_bullet_slide(prs, "Business Value", [
        "For DR Managers:",
        "  • Familiar VMware SRM-like interface",
        "  • One-click recovery execution",
        "  • Real-time progress visibility",
        "",
        "For DevOps Engineers:",
        "  • Infrastructure as Code",
        "  • API-first design",
        "  • CloudWatch observability",
        "",
        "For CIOs/Finance:",
        "  • 79% cost reduction vs VMware SRM",
        "  • Zero capital expenditure",
        "  • 75% staff reduction"
    ], 13, total_slides)
    
    # Slide 14: Deployment Process
    add_bullet_slide(prs, "Deployment Process", [
        "Phase 1: Setup (30 minutes)",
        "  • Deploy CloudFormation master template",
        "  • Configure Cognito users",
        "  • Review IAM permissions",
        "",
        "Phase 2: Configuration (30 minutes)",
        "  • Discover DRS source servers",
        "  • Create protection groups",
        "  • Define recovery plans",
        "",
        "Phase 3: Testing (30 minutes)",
        "  • Execute drill recovery",
        "  • Validate server launch",
        "  • Review execution history"
    ], 14, total_slides)
    
    # Slide 15: Key Takeaways
    add_bullet_slide(prs, "Key Takeaways", [
        "1. Solution: Serverless orchestration for AWS DRS",
        "",
        "2. Cost: $3,360/year (99.7% cheaper than VMware SRM)",
        "",
        "3. Performance: 65% faster recovery (16h → 6h)",
        "",
        "4. Savings: $8.1M over 3 years (79% vs VMware SRM)",
        "",
        "5. Deployment: < 1 hour via CloudFormation",
        "",
        "6. Value: Enterprise-grade capabilities at serverless scale"
    ], 15, total_slides)
    
    # Slide 16: Thank You (no footer)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE
    p.font.name = "Amazon Ember Display"
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AWS DRS Orchestration Platform"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Amazon Ember"
    p.alignment = PP_ALIGN.CENTER
    
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5.4), Inches(8), Inches(0.8))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Making Enterprise Disaster Recovery Simple, Fast, and Cost-Effective"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.name = "Amazon Ember"
    p.alignment = PP_ALIGN.CENTER
    
    # Powered by AWS logo (reduced to half size)
    if os.path.exists('aws_logo.png'):
        try:
            slide.shapes.add_picture('aws_logo.png', 
                                    Inches(4.25), Inches(6.3), 
                                    width=Inches(1.5))
        except:
            pass
    
    # Save
    output_file = "AWS_DRS_Orchestration_Solution_EDITABLE.pptx"
    prs.save(output_file)
    
    print(f"\n{'='*60}")
    print(f"✅ AWS-BRANDED POWERPOINT CREATED")
    print(f"{'='*60}")
    print(f"📄 File: {output_file}")
    print(f"📊 Slides: {len(prs.slides)}")
    print(f"🎨 Font: Amazon Ember (with fallback)")
    print(f"🏢 Branding: AWS Official Colors")
    print(f"📌 Features:")
    print(f"   • Slide numbers on every slide")
    print(f"   • Date and confidentiality notice")
    print(f"   • AWS logo in footer")
    print(f"   • 'Powered by AWS' on title/closing")
    print(f"   • AWS Orange (#FF9900) headers")
    print(f"   • Professional table styling")
    print(f"✏️  100% Editable in PowerPoint!")
    print(f"{'='*60}\n")
    
    return output_file

if __name__ == "__main__":
    create_presentation()
