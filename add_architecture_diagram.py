#!/usr/bin/env python3
"""
Add Architecture Diagram to PowerPoint Presentation

This script adds the architecture diagram image to an existing PowerPoint presentation,
creating a new slide with the diagram properly sized and positioned.
"""

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import sys
import os

def add_architecture_slide(pptx_path, image_path, output_path=None, insert_position=4):
    """
    Add architecture diagram to PowerPoint presentation
    
    Args:
        pptx_path: Path to existing PowerPoint file
        image_path: Path to architecture diagram image
        output_path: Optional output path (defaults to overwriting input)
        insert_position: Position to insert slide (default: 4, between slides 4 and 5)
    """
    
    # Load existing presentation
    prs = Presentation(pptx_path)
    
    # Check if architecture slide already exists
    architecture_slide_index = None
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title and "Architecture" in slide.shapes.title.text:
            architecture_slide_index = i
            print(f"ℹ️  Found existing architecture slide at position {i+1}")
            break
    
    if architecture_slide_index is None:
        # Add new slide at specified position
        # Use blank layout (index 6 is typically blank)
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Move slide to correct position
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])  # Remove from end
        xml_slides.insert(insert_position, slides[-1])  # Insert at position
        
        print(f"✅ New slide created and inserted at position {insert_position + 1}")
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.75)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "Solution Architecture"
        
        # Format title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Inches(0.5)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(255, 153, 0)  # Amazon orange
    else:
        slide = prs.slides[architecture_slide_index]
    
    # Add image to slide
    # Position: centered, with margins
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)  # Let height auto-scale
    
    try:
        slide.shapes.add_picture(image_path, left, top, width=width)
        print(f"✅ Architecture diagram added successfully!")
    except Exception as e:
        print(f"❌ Error adding image: {e}")
        return False
    
    # Save presentation
    if output_path is None:
        output_path = pptx_path
    
    prs.save(output_path)
    print(f"✅ PowerPoint saved: {output_path}")
    return True

def main():
    """Main function"""
    
    # Default paths
    pptx_path = "AWS_DRS_Orchestration_Solution_EDITABLE.pptx"
    image_path = "docs/architecture-diagram.png"
    
    # Check if files exist
    if not os.path.exists(pptx_path):
        print(f"❌ Error: PowerPoint file not found: {pptx_path}")
        sys.exit(1)
    
    if not os.path.exists(image_path):
        print(f"❌ Error: Architecture diagram not found: {image_path}")
        sys.exit(1)
    
    print("📊 Adding Architecture Diagram to PowerPoint")
    print("=" * 60)
    print(f"📄 PowerPoint: {pptx_path}")
    print(f"🖼️  Diagram: {image_path}")
    print()
    
    # Add diagram between slides 4 and 5 (index 4)
    success = add_architecture_slide(pptx_path, image_path, insert_position=4)
    
    if success:
        print()
        print("✅ Operation Complete!")
        print(f"📁 Updated file: {pptx_path}")
        print()
        print("📋 Next Steps:")
        print("1. Open the PowerPoint file")
        print("2. Navigate to 'Solution Architecture' slide")
        print("3. Verify diagram appears correctly")
        print("4. Adjust size/position if needed")
    else:
        print()
        print("❌ Operation Failed")
        print("Please check error messages above")
        sys.exit(1)

if __name__ == "__main__":
    main()
