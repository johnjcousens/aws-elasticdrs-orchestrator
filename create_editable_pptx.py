#!/usr/bin/env python3
"""
Create native editable PowerPoint from AWS DRS Orchestration presentation.
All text will be fully editable in PowerPoint.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# AWS Brand Colors
AWS_ORANGE = RGBColor(255, 153, 0)  # #FF9900
AWS_DARK = RGBColor(35, 47, 62)     # #232F3E

def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "AWS DRS Orchestration Platform"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = AWS_ORANGE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    tf.text = "Serverless Disaster Recovery Orchestration for AWS Elastic Disaster Recovery"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    tf = tag_box.text_frame
    tf.text = "Enterprise-Grade Solution for Automated Recovery"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_bullet_slide(prs, title, bullets):
    """Add slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    # Title
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = AWS_DARK
    
    # Content
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)

def add_image_slide(prs, title, image_path):
    """Add slide with image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = AWS_DARK
    
    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = subtitle_box.text_frame
    tf.text = "100% Serverless | Zero Infrastructure | Pay-per-Use"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_table_slide(prs, title, headers, rows):
    """Add slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = AWS_DARK
    
    # Table
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5)
    ).table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = AWS_ORANGE
    
    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(prs)
    
    # Slide 2: What is AWS DRS Orchestration?
    add_bullet_slide(prs, "What is AWS DRS Orchestration?", [
        "Wave-based recovery execution (like VMware SRM)",
        "Server dependency management",
        "One-click drill testing",
        "Real-time progress monitoring",
        "Complete audit trail",
        "",
        "Built with: 100% AWS serverless technologies",
        "Deployment: < 1 hour via CloudFormation",
        "Cost: $3,360/year (fixed, regardless of server count)"
    ])
    
    # Slide 3: The Problem
    add_bullet_slide(prs, "The Problem We Solve", [
        "WITHOUT ORCHESTRATION:",
        "  • Manual server recovery (error-prone)",
        "  • No dependency management",
        "  • 12-16 hour recovery times",
        "  • Complex drill testing procedures",
        "  • No visibility into recovery progress",
        "",
        "MIGRATION CHALLENGE:",
        "  • Customers have VMware SRM orchestration",
        "  • AWS DRS lacks equivalent capabilities",
        "  • This gap blocks AWS DRS adoption"
    ])
    
    # Slide 4: Architecture Diagram
    add_image_slide(prs, "The Solution Architecture", "docs/architecture-diagram.png")
    
    # Slide 5: Architecture Components
    add_table_slide(prs, "Architecture Components", 
        ["Layer", "Technology", "Purpose"],
        [
            ["Frontend", "React + TypeScript + MUI", "Modern web UI"],
            ["CDN", "CloudFront + S3", "Global distribution"],
            ["Authentication", "Cognito User Pools", "Secure access"],
            ["API", "API Gateway (REST)", "Managed API layer"],
            ["Compute", "Lambda (Python 3.12)", "Business logic"],
            ["Orchestration", "DRS API Integration", "Recovery execution"],
            ["Storage", "DynamoDB (3 tables)", "Configuration data"],
            ["Monitoring", "CloudWatch + CloudTrail", "Logs and audit"]
        ]
    )
    
    # Slide 6: Core Capabilities
    add_bullet_slide(prs, "Core Capabilities", [
        "1. Protection Groups",
        "  • Group servers logically",
        "  • Auto-discover DRS servers",
        "  • Reusable across plans",
        "",
        "2. Recovery Plans",
        "  • Multi-wave recovery sequences",
        "  • Wave 1: Database | Wave 2: Application | Wave 3: Web",
        "  • Dependencies between waves"
    ])
    
    # Slide 7: Wave-Based Execution
    add_bullet_slide(prs, "Wave-Based Execution", [
        "3. Wave-Based Execution",
        "  • 15-second delays between servers",
        "  • 30-second delays between waves",
        "  • Automatic retry with exponential backoff",
        "",
        "4. Drill Mode Testing",
        "  • Creates test instances (not production cutover)",
        "  • Zero impact on replication"
    ])
    
    # Slide 8: Monitoring & Recovery
    add_bullet_slide(prs, "Real-Time Monitoring & One-Click Recovery", [
        "5. Real-Time Monitoring",
        "  • Execution history (start/end times)",
        "  • Job status tracking",
        "  • CloudWatch integration",
        "  • Audit trail in CloudTrail",
        "",
        "6. One-Click Recovery",
        "  • Select recovery plan",
        "  • Choose drill vs production mode",
        "  • Click 'Execute Recovery'",
        "  • Monitor progress in real-time"
    ])
    
    # Slide 9: Performance Metrics
    add_table_slide(prs, "Performance Metrics",
        ["Metric", "Before (Manual)", "After (Orchestrated)", "Improvement"],
        [
            ["Recovery Time", "12-16 hours", "4-6 hours", "65% faster"],
            ["Setup Time", "2-3 weeks", "< 1 hour", "99% faster"],
            ["Drill Execution", "4-5 hours", "2-3 minutes", "98% faster"],
            ["Error Rate", "15-20%", "< 2%", "90% reduction"]
        ]
    )
    
    # Slide 10: Cost Breakdown
    add_table_slide(prs, "AWS DRS Orchestration Platform Cost",
        ["Service", "Monthly", "Annual", "Notes"],
        [
            ["Lambda", "$100", "$1,200", "Recovery execution"],
            ["DynamoDB", "$50", "$600", "3 tables (on-demand)"],
            ["API Gateway", "$50", "$600", "REST API requests"],
            ["CloudFront", "$20", "$240", "Global CDN"],
            ["S3", "$10", "$120", "Frontend hosting"],
            ["Cognito", "$20", "$240", "User authentication"],
            ["CloudWatch", "$20", "$240", "Log retention"],
            ["CloudTrail", "$10", "$120", "Audit logging"],
            ["TOTAL", "$280", "$3,360", "Fixed cost"]
        ]
    )
    
    # Slide 11: Total Solution Cost
    add_table_slide(prs, "Total Solution Cost (1,000 Servers)",
        ["Component", "Annual Cost", "Notes"],
        [
            ["AWS DRS Service", "$480,000", "Scales with server count"],
            ["+ Orchestration Platform", "$3,360", "Fixed serverless cost"],
            ["= Total Solution", "$483,360", "Complete DR solution"]
        ]
    )
    
    # Slide 12: Competitive Comparison
    add_table_slide(prs, "Competitive Comparison",
        ["Solution", "Orchestration Cost", "Setup Time", "Serverless", "Wave Support"],
        [
            ["AWS DRS Orchestration", "$3.4K/year", "< 1 hour", "✅ Yes", "✅ Yes"],
            ["VMware SRM", "$1M+/year", "Weeks", "❌ No", "✅ Yes"],
            ["Zerto", "$100K+/year", "Days", "❌ No", "⚠️ Limited"],
            ["Azure Site Recovery", "$50K+/year", "Days", "❌ No", "❌ No"]
        ]
    )
    
    # Slide 13: TCO Summary
    add_table_slide(prs, "TCO Summary (3-Year, 1,000 Servers)",
        ["Solution", "3-Year Total", "Annual Avg", "Staffing"],
        [
            ["VMware SRM", "$10.2M", "$3.4M", "4 FTEs"],
            ["AWS DRS + Orch", "$2.1M", "$697K", "1 FTE"],
            ["Savings", "$8.1M (79%)", "$2.7M", "3 FTEs"]
        ]
    )
    
    # Slide 14: Business Value
    add_bullet_slide(prs, "Business Value", [
        "For DR Managers:",
        "  • Familiar VMware SRM-like interface",
        "  • One-click recovery execution",
        "  • Real-time progress visibility",
        "",
        "For DevOps Engineers:",
        "  • Infrastructure as Code",
        "  • API-first design",
        "  • CloudWatch observability",
        "",
        "For CIOs/Finance:",
        "  • 79% cost reduction vs VMware SRM",
        "  • Zero capital expenditure",
        "  • 75% staff reduction"
    ])
    
    # Slide 15: Key Takeaways
    add_bullet_slide(prs, "Key Takeaways", [
        "1. Solution: Serverless orchestration for AWS DRS",
        "",
        "2. Cost: $3,360/year (99.7% cheaper than VMware SRM)",
        "",
        "3. Performance: 65% faster recovery (16h → 6h)",
        "",
        "4. Savings: $8.1M over 3 years (79% vs VMware SRM)",
        "",
        "5. Deployment: < 1 hour via CloudFormation",
        "",
        "6. Value: Enterprise-grade capabilities at serverless scale"
    ])
    
    # Slide 16: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AWS DRS Orchestration Platform"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Making Enterprise Disaster Recovery Simple, Fast, and Cost-Effective"
    p.font.size = Pt(18)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    output_file = "AWS_DRS_Orchestration_Solution_EDITABLE.pptx"
    prs.save(output_file)
    print(f"✅ Created editable PowerPoint: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"✏️  All text is fully editable!")
    return output_file

if __name__ == "__main__":
    create_presentation()
