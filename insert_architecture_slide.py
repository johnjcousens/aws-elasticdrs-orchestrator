#!/usr/bin/env python3
"""
Insert Architecture Slide into PowerPoint Presentation (Safe Method)

This script safely inserts the architecture diagram by creating a new presentation
and copying slides in the correct order, avoiding XML manipulation that can corrupt files.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys
import os

def copy_slide(source_prs, dest_prs, slide_index):
    """
    Copy a slide from source to destination presentation
    Copies layout and all shapes
    """
    source_slide = source_prs.slides[slide_index]
    
    # Get the slide layout (use blank if not available)
    try:
        slide_layout = dest_prs.slide_layouts[source_slide.slide_layout.slide_layout_id]
    except:
        slide_layout = dest_prs.slide_layouts[6]  # Blank layout
    
    # Add new slide
    dest_slide = dest_prs.slides.add_slide(slide_layout)
    
    # Copy all shapes from source slide
    for shape in source_slide.shapes:
        # Get shape properties
        el = shape.element
        newel = el.__class__(el)
        dest_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return dest_slide

def create_architecture_slide(prs, image_path):
    """
    Create architecture diagram slide
    """
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.75)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Solution Architecture"
    
    # Format title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(40)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 153, 0)  # Amazon orange
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Add image
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    
    try:
        slide.shapes.add_picture(image_path, left, top, width=width)
        print(f"✅ Architecture diagram added successfully!")
        return True
    except Exception as e:
        print(f"❌ Error adding image: {e}")
        return False

def insert_architecture_slide_safely(source_pptx, image_path, output_pptx, insert_after_slide=4):
    """
    Safely insert architecture slide by rebuilding presentation
    
    Args:
        source_pptx: Source PowerPoint file
        image_path: Path to architecture diagram
        output_pptx: Output PowerPoint file
        insert_after_slide: Insert after this slide number (1-indexed)
    """
    print(f"📖 Reading source presentation: {source_pptx}")
    source_prs = Presentation(source_pptx)
    total_slides = len(source_prs.slides)
    print(f"   Found {total_slides} slides")
    
    print(f"🏗️  Creating new presentation...")
    dest_prs = Presentation()
    
    # Remove default slide
    rId = dest_prs.slides._sldIdLst[0].rId
    dest_prs.part.drop_rel(rId)
    del dest_prs.slides._sldIdLst[0]
    
    # Copy slides before insertion point
    print(f"📄 Copying slides 1-{insert_after_slide}...")
    for i in range(insert_after_slide):
        copy_slide(source_prs, dest_prs, i)
    
    # Add architecture slide
    print(f"🖼️  Inserting architecture slide at position {insert_after_slide + 1}...")
    if not create_architecture_slide(dest_prs, image_path):
        return False
    
    # Copy remaining slides
    if insert_after_slide < total_slides:
        print(f"📄 Copying slides {insert_after_slide + 1}-{total_slides}...")
        for i in range(insert_after_slide, total_slides):
            copy_slide(source_prs, dest_prs, i)
    
    # Save
    print(f"💾 Saving to: {output_pptx}")
    dest_prs.save(output_pptx)
    print(f"✅ PowerPoint saved successfully!")
    return True

def main():
    """Main function"""
    
    # Paths
    source_pptx = "AWS_DRS_Orchestration_Solution_EDITABLE.pptx"
    image_path = "docs/architecture-diagram.png"
    output_pptx = "AWS_DRS_Orchestration_Solution_EDITABLE.pptx"
    
    # Check files
    if not os.path.exists(source_pptx):
        print(f"❌ Error: Source PowerPoint not found: {source_pptx}")
        sys.exit(1)
    
    if not os.path.exists(image_path):
        print(f"❌ Error: Architecture diagram not found: {image_path}")
        sys.exit(1)
    
    print("📊 Inserting Architecture Diagram into PowerPoint")
    print("=" * 60)
    print(f"📄 Source: {source_pptx}")
    print(f"🖼️  Diagram: {image_path}")
    print(f"📍 Insert Position: After slide 4")
    print()
    
    # Create temporary output first
    temp_output = "temp_presentation.pptx"
    
    try:
        # Insert slide safely
        success = insert_architecture_slide_safely(
            source_pptx, 
            image_path, 
            temp_output,
            insert_after_slide=4
        )
        
        if success:
            # Replace original with new version
            os.replace(temp_output, output_pptx)
            print()
            print("✅ Operation Complete!")
            print(f"📁 Updated file: {output_pptx}")
            print()
            print("📋 Slide Order:")
            print("   1. Title Slide")
            print("   2. Executive Summary")
            print("   3. Problem Statement")
            print("   4. Solution Overview")
            print("   5. Solution Architecture ← NEW (Architecture Diagram)")
            print("   6. Key Features")
            print("   7-17. Additional slides")
        else:
            print()
            print("❌ Operation Failed")
            if os.path.exists(temp_output):
                os.remove(temp_output)
            sys.exit(1)
            
    except Exception as e:
        print(f"❌ Error: {e}")
        if os.path.exists(temp_output):
            os.remove(temp_output)
        sys.exit(1)

if __name__ == "__main__":
    main()
