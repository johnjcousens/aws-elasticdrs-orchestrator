#!/usr/bin/env python3
"""
Add Architecture Slide to PowerPoint (Simple & Safe)

This script adds an architecture slide to the end of the presentation.
Users can manually move it to the desired position in PowerPoint.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys
import os

def add_architecture_slide(pptx_path, image_path):
    """Add architecture slide to presentation"""
    
    print(f"📖 Loading presentation: {pptx_path}")
    prs = Presentation(pptx_path)
    print(f"   Current slides: {len(prs.slides)}")
    
    # Use blank layout
    print(f"🏗️  Creating architecture slide...")
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.75)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Solution Architecture"
    
    # Format title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(40)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 153, 0)  # Amazon orange
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Add image
    print(f"🖼️  Adding architecture diagram...")
    img_left = Inches(0.5)
    img_top = Inches(1.5)
    img_width = Inches(9)
    
    try:
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
        print(f"✅ Architecture diagram added successfully!")
    except Exception as e:
        print(f"❌ Error adding image: {e}")
        return False
    
    # Save
    print(f"💾 Saving presentation...")
    prs.save(pptx_path)
    print(f"✅ Saved: {pptx_path}")
    print(f"   Total slides: {len(prs.slides)}")
    
    return True

def main():
    """Main function"""
    
    # Paths
    pptx_path = "AWS_DRS_Orchestration_Solution_EDITABLE.pptx"
    image_path = "docs/architecture-diagram.png"
    
    # Check files
    if not os.path.exists(pptx_path):
        print(f"❌ Error: PowerPoint not found: {pptx_path}")
        sys.exit(1)
    
    if not os.path.exists(image_path):
        print(f"❌ Error: Architecture diagram not found: {image_path}")
        sys.exit(1)
    
    print("📊 Adding Architecture Slide to PowerPoint")
    print("=" * 60)
    print(f"📄 PowerPoint: {pptx_path}")
    print(f"🖼️  Diagram: {image_path}")
    print()
    
    # Add slide
    success = add_architecture_slide(pptx_path, image_path)
    
    if success:
        print()
        print("✅ Operation Complete!")
        print()
        print("📋 Next Steps:")
        print("1. Open PowerPoint file")
        print("2. Find 'Solution Architecture' slide (added at end)")
        print("3. Drag slide to position 5 (between slides 4 and 5)")
        print("4. Save the file")
        print()
        print("💡 Tip: In PowerPoint, you can drag slides in the")
        print("   left sidebar to reorder them")
    else:
        print()
        print("❌ Operation Failed")
        sys.exit(1)

if __name__ == "__main__":
    main()
